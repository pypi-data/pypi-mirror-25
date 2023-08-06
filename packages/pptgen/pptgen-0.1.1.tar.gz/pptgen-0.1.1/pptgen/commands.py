"""Python-PPTX customized module."""
import os
import six
import copy
import logging
import requests
import datetime
import tempfile
import collections
import numpy as np
from . import utils
import pandas as pd
import matplotlib.cm
from . import fontwidth
import matplotlib.colors
from . import color as _color
from pptx.dml.color import RGBColor
from tornado.template import Template
from tornado.escape import to_unicode
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import XyChartData
from pptx.chart.data import BubbleChartData
from six.moves.urllib_parse import urlparse
from gramex.transforms import build_transform

_template_cache = {}


def template(tmpl, data):
    """Execute tornado template."""
    if tmpl not in _template_cache:
        _template_cache[tmpl] = Template(tmpl, autoescape=None)
    return to_unicode(_template_cache[tmpl].generate(**data))


def text(shape, spec, data):
    '''
    Replace entire text of shape with spec['text'].
    '''
    if not shape.has_text_frame:
        logging.error('"%s" is not a TextShape to apply text:', shape.name)
        return

    run_flag = True
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = template(spec['text'], data) if run_flag else ''
            run_flag = False


def replace(shape, spec, data):
    '''
    Replace keywords in shape using the dictionary at spec['replace'].
    '''
    if not shape.has_text_frame:
        logging.error('"%s" is not a TextShape to apply text:', shape.name)
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old, new in spec['replace'].items():
                run.text = run.text.replace(old, template(new, data))


def image(shape, spec, data):
    '''
    Replace image with a different file specified in spec['image']
    '''
    image = template(spec['image'], data)
    # If it's a URL, use the requests library's raw stream as a file-like object
    if urlparse(image).netloc:
        r = requests.get(image)
        with tempfile.NamedTemporaryFile(delete=False) as handle:
            handle.write(r.content)
        new_img_part, new_rid = shape.part.get_or_add_image_part(handle.name)
        os.unlink(handle.name)
    else:
        new_img_part, new_rid = shape.part.get_or_add_image_part(image)
    old_rid = shape._pic.blip_rId
    shape._pic.blipFill.blip.rEmbed = new_rid
    shape.part.related_parts[old_rid].blob = new_img_part.blob


def stack_elements(replica, shape, stack=False, margin=None):
    '''
    Function to extend elements horizontally or vertically.
    '''
    config = {'vertical': {'axis': 'y', 'attr': 'height'},
              'horizontal': {'axis': 'x', 'attr': 'width'}}
    if stack is not None:
        grp_sp = shape.element
        # Adding a 15% margin between original and new object.
        margin = 0.50 if not margin else margin
        for index in range(replica - 1):
            # Adding a cloned object to shape
            extend_shape = copy.deepcopy(grp_sp)
            # Getting attributes and axis values from config based on stack.
            attr = config.get(stack, {}).get('attr', 0)
            axis = config.get(stack, {}).get('axis', 0)
            # Taking width or height based on stack value and setting a margin.
            metric_val = getattr(shape, attr)
            axis_val = getattr(extend_shape, axis)
            # Setting margin accordingly either vertically or horizontally.
            axis_pos = metric_val * (index + 1)
            set_attr = axis_val + axis_pos + (axis_pos * margin)
            # Setting graphic position of newly created object to slide.
            setattr(extend_shape, axis, int(set_attr))
            # Adding newly created object to slide.
            grp_sp.addnext(extend_shape)


def stack_shapes(collection, change, data):
    '''
    Function to stack Shapes if required.
    '''
    shape_data = data
    for shape in collection:
        if shape.name not in change:
            continue
        info = change[shape.name]
        if 'data' in info and info.get('stack') is not None:
            shape_data = eval(info.get('data'))
        stack_elements(len(shape_data), shape, stack=info.get('stack'),
                       margin=info.get('margin'))


def generate_slide(prs, source):
    '''
    Create a slide layout.
    '''
    layout_items_count = [
        len(layout.placeholders) for layout in prs.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return prs.slide_layouts[blank_layout_id]


def copy_slide_elem(shape, dest):
    '''
    Function to copy slide elements into a newly created slide.
    '''
    if dest:
        el = shape.element
        new_elem = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(new_elem, 'p:extLst')


def add_new_slide(dest, source_slide):
    '''
    Function to add a new slide to presentation.
    '''
    if dest:
        for key, value in six.iteritems(source_slide.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                dest.part.rels.add_relationship(value.reltype, value._target, value.rId)


def replicate_slide(change, prs, data):
    '''
    Add Duplicate slide to the presentation.
    '''
    if 'replicate-slides' in change and change['replicate-slides'] is not None:
        replica_conf = change['replicate-slides']['slide-number']
        for conf in replica_conf:
            for slide_number, replica_logic in conf.items():
                replica = template(replica_logic, data)
                for replica in range(int(replica) - 1):
                    source = prs.slides[slide_number - 1]
                    add_new_slide(prs, source)


def move_slide(presentation, old_index, new_index):
    '''
    Move a slide's index number.
    '''
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def delete_slide(presentation, index):
    '''
    Delete a slide from Presentation.
    '''
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def rect_css(shape, **kwargs):
    """Function to add text to shape."""
    for key in {'fill', 'stroke'}:
        if kwargs.get(key):
            fill = shape.fill if key == 'fill' else shape.line.fill
            rectcss = kwargs[key].rsplit('#')[-1].lower()
            rectcss = rectcss + ('0' * (6 - len(rectcss)))
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(rectcss)


def add_text_to_shape(shape, text, font_size, txt_fill):
    """Function to add text to shape."""
    min_inc = 13000
    pixel_inch = 10000
    font_size = max(font_size * pixel_inch, min_inc)
    if font_size > min_inc:
        txt_fill = txt_fill.rsplit('#')[-1].lower()
        txt_fill = txt_fill + (txt_fill[0] * (6 - len(txt_fill)))
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.add_run()
        for run in paragraph.runs:
            run.text = text
            shape_txt = run.font
            shape_txt = run.font.fill
            shape_txt.solid()
            run.font.size = '{:.0f}'.format(font_size)
            run.font.color.rgb = RGBColor.from_string(txt_fill)


def scale_data(data, lo, hi, factor):
    """Function to scale data."""
    return ((data - lo) / (hi - lo)) * factor


def rect(shape, x, y, width, height):
    """Add rectangle to slide."""
    return shape.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)


def _prepare_data(config, data):
    '''
    Prepare chart data.
    '''
    if not isinstance(data, (list, pd.DataFrame)):
        raise NotImplementedError()

    if isinstance(data, list):
        data = pd.DataFrame(data)

    # Filtering data based on query
    if config.get('query') is not None and 'query' in config:
        data = data.query(config['query'])

    # Create a new dataset with the relevant columns
    if config.get('columns') is not None and 'columns' in config:
        for name, column_info in config['columns'].items():
            if 'value' in column_info:
                data[name] = data[column_info['value']]
            elif 'formula' in column_info:
                data[name] = data.eval(column_info['formula'])
    return data


def _update_chart(info, data, chart_data, series_columns, chart='line'):
    '''
    Updating XyChart data.
    '''
    series_dict = {}
    columns = list(data.columns)
    if info['x'] in columns:
        columns.remove(info['x'])

    if chart == 'line':
        for series in series_columns:
            if np.issubdtype(data[series].dtype, np.number):
                chart_data.add_series(
                    series, tuple(data[series].dropna().tolist()))
        return chart_data

    for index, row in data.fillna(0).iterrows():
        for col in series_columns:
            if col in columns and np.issubdtype(data[col].dtype, np.number):
                serieslist = [series.name for series in chart_data._series]
                if col not in serieslist:
                    series_dict[col] = chart_data.add_series(col)
                if chart == 'scatter':
                    try:
                        x = float(row[info['x']])
                    except TypeError:
                        x = index
                    series_dict[col].add_data_point(x, row[col])
                elif chart == 'bubble':
                    bubble_size = 1
                    if info.get('size') is not None and 'size' in info:
                        bubble_size = row[info['size']]
                    series_dict[col].add_data_point(
                        row[info['x']], row[col], bubble_size)
    return chart_data


def _extract_tbl_properties(table, tbl_style, row_idx, rowtype):
    '''
    Exctracting table properties.
    '''
    # Get Table rows/columns properties.
    row_text_tag = ['./a:p/a:endParaRPr/a:solidFill/a:srgbClr',
                    './a:p/a:endParaRPr/a:solidFill/a:schemeClr']
    txt_clr_map = {
        './a:p/a:endParaRPr/a:solidFill/a:schemeClr': None,
        './a:p/a:endParaRPr/a:solidFill/a:srgbClr': None
    }

    text_dict = {}
    for row_text in row_text_tag:
        row_txt_clr = table.rows[row_idx].cells[0]._tc.txBody.xpath(row_text)
        clr_key = '{}_{}'.format(rowtype, row_text.rsplit(':')[-1])
        if len(row_txt_clr) > 0:
            text_dict[clr_key] = row_txt_clr[0].val
        else:
            text_dict[clr_key] = txt_clr_map.get(row_text)

    # Text inside Cell properties information.
    info_txt = table.rows[row_idx].cells[0].text_frame.paragraphs[0]
    if not hasattr(info_txt, 'runs'):
        info_txt.add_run()
    if info_txt.runs:
        info_txt = info_txt.runs[0].font
        text_dict['{}{}'.format(rowtype, '_italic')] = info_txt.italic
        text_dict['{}{}'.format(rowtype, '_bold')] = info_txt.bold
        text_dict['{}{}'.format(rowtype, '_txt_size')] = info_txt.size
        text_dict['{}{}'.format(rowtype, '_font_name')] = info_txt.name
        text_dict['{}{}'.format(rowtype, '_underline')] = info_txt.underline

    # Cell properties information.
    cell_info = ['./a:srgbClr', './a:schemeClr']
    cell_info_map = {'./a:schemeClr': None, './a:srgbClr': None}

    cell_dict = {}
    info_cell = table.rows[row_idx].cells[0]._tc.tcPr.solidFill
    if info_cell is not None:
        for cell_prop in cell_info:
            cell_clr = info_cell.xpath(cell_prop)
            cell_key = '{}_{}'.format(rowtype, cell_prop.rsplit(':')[-1])
            if len(cell_clr) > 0:
                cell_dict[cell_key] = cell_clr[0].val
            else:
                cell_dict[cell_key] = cell_info_map.get(cell_prop)
    tbl_style['text'] = text_dict
    tbl_style['cell'] = cell_dict
    return tbl_style


def _extend_table(shape, data, total_rows, total_columns):
    '''
    Function to extend table rows and columns if required.
    '''
    avail_rows = len(shape.table.rows)
    avail_cols = len(shape.table.columns)

    col_width = shape.table.columns[0].width
    row_height = shape.table.rows[0].height
    # Extending Table Rows if required based on the data
    while avail_rows < total_rows:
        shape.table.rows._tbl.add_tr(row_height)
        avail_rows += 1
    # Extending Table Columns if required based on the data
    while avail_cols < total_columns:
        shape.table._tbl.tblGrid.add_gridCol(col_width)
        avail_cols += 1


def _table_text_css(run, table_style, rowtype, text):
    '''
    Function to Apply table styles for rows and columns.
    '''
    run.text = '{}'.format(text)
    rows_text = run.font.fill
    rows_text.solid()
    rgb = None
    theme_color = None
    if 'srgbClr' in table_style:
        rgb = table_style.get('{}_{}'.format(rowtype, 'srgbClr'))
    if 'schemeClr' in table_style:
        theme_color = table_style.get('{}_{}'.format(rowtype, 'schemeClr'))

    if theme_color:
        rows_text.fore_color.theme_color = theme_color
    else:
        if rgb:
            run.font.color.rgb = RGBColor.from_string(rgb) if isinstance(rgb, str) else rgb
    if '{}{}'.format(rowtype, '_font_name') in table_style:
        run.font.name = table_style['{}{}'.format(rowtype, '_font_name')]
    if '{}{}'.format(rowtype, '_txt_size') in table_style:
        run.font.size = table_style['{}{}'.format(rowtype, '_txt_size')]
    if '{}{}'.format(rowtype, '_bold') in table_style:
        run.font.bold = table_style['{}{}'.format(rowtype, '_bold')]
    if '{}{}'.format(rowtype, '_italic') in table_style:
        run.font.italic = table_style['{}{}'.format(rowtype, '_italic')]
    if '{}{}'.format(rowtype, '_underline') in table_style:
        run.font.underline = table_style['{}{}'.format(rowtype, '_underline')]


def _table_cell_css(shape, cell, cell_prop, row_type, gradient):
    '''
    Set cells properties.
    '''
    props = cell.fill
    props.solid()

    if gradient:
        if isinstance(gradient, str):
            clr = gradient.replace('#', '') + ('0' * (6 - len(gradient.replace('#', ''))))
            props.fore_color.rgb = RGBColor.from_string(clr)
        else:
            props.fore_color.rgb = gradient
        return

    theme_color = cell_prop.get('{}_{}'.format(row_type, 'schemeClr'))
    if theme_color:
        props.fore_color.theme_color = theme_color
        # Setting brightness - Not sure when it gets failed
        bright = shape.cells[0]._tc.tcPr
        if bright.solidFill.schemeClr.lumMod:
            props.fore_color.brightness = bright.solidFill.schemeClr.lumMod.val - 1.0
        elif bright.solidFill.schemeClr.lumOff:
            props.fore_color.brightness = bright.solidFill.schemeClr.lumOff.val
    else:
        rgb = cell_prop.get('{}_{}'.format(row_type, 'srgbClr')) or 'ffffff'
        if rgb:
            props.fore_color.rgb = RGBColor.from_string(rgb) if isinstance(rgb, str) else rgb


def _table_css(shape, rowtype='header'):
    '''
    Function to get Table style for rows and columns.
    '''
    tbl_style = {}
    table = shape.table
    row_idx = 0 if rowtype == 'header' else 1
    # row_idx = 1 if len(table.rows) > 1 else 0
    tbl_style = _extract_tbl_properties(table, tbl_style, row_idx, rowtype)
    # Get Table Header style' from here yet to implement.
    # cell_clr = table.rows[0].cells[0]._tc.tcPr.solidFill.schemeClr.val
    # tbl_style['margin_bottom'] = cell_txt_clr.margin_bottom
    return tbl_style


def _table_cell_gradient(info, data):
    '''
    Function to calculate color gradient.
    '''
    color_grad = None
    if info.get('gradient') and 'gradient' in info:
        get_grad = info.get('gradient')
        color_grad = _color.RdYlGn
        data = pd.DataFrame(data) if not isinstance(data, (pd.DataFrame,)) else data
        numeric = data._get_numeric_data()
        min_data = min([numeric[x].min() for x in numeric.columns])
        max_data = max([numeric[x].max() for x in numeric.columns])
        mean_data = (min_data + max_data) / 2.0
        grad_map = {0: min_data, 1: mean_data, 2: max_data}

        if isinstance(get_grad, str):
            color_grad = getattr(_color, get_grad)

        elif isinstance(get_grad, list) and isinstance(get_grad[0], tuple):
            color_grad = get_grad

        elif isinstance(get_grad, list) and isinstance(get_grad[0], str):
            color_grad = [(grad_map[index], clr) for index, clr in enumerate(get_grad)]
    return color_grad


def compile_function(spec, key, data, handler):
    """A function to compile configuration."""
    if not spec.get(key):
        return None
    _vars = {'_color': None, 'data': None, 'handler': None}
    if not isinstance(spec[key], (dict,)):
        spec[key] = {'function': '{}'.format(spec[key])}
    elif isinstance(spec[key], (dict,)) and 'function' not in spec[key]:
        spec[key] = {'function': '{}'.format(spec[key])}
    args = {'data': data, 'handler': handler, '_color': _color}
    return build_transform(spec[key], vars=_vars)(**args)[0]


def table(shape, spec, data):
    """Update an existing Table shape with data."""
    if not spec.get('table', {}).get('data'):
        return

    spec = spec['table']
    handler = data.pop('handler') if 'handler' in data else None
    data = compile_function(spec, 'data', data, handler)
    color_grad = _table_cell_gradient(spec, data)

    if isinstance(data, (pd.DataFrame,)):
        data = data.to_dict(orient='records')
    if not len(data):
        return

    data_cols = list(data[0].keys())
    # Extending table if required.
    _extend_table(shape, data, len(data) + 1, len(data_cols))
    # Fetching Table Style for All Cells and texts.
    for row_num, row in enumerate(shape.table.rows):
        cols = len(row.cells._tr.tc_lst)
        # Extending cells in newly added rows.
        while cols < len(data_cols):
            row.cells._tr.add_tc()
            cols += 1
        # row_type = 'header' if row_num == 0 else 'row'
        row_type = 'header'
        tbl_style = _table_css(shape, rowtype='header')
        for col_num, cell in enumerate(row.cells):
            colname = data_cols[col_num]
            for curr_cell in cell.text_frame.paragraphs:
                if not curr_cell.text.strip():
                    curr_cell.add_run()
                for run in curr_cell.runs:
                    text_prop = tbl_style.get('text', {})
                    txt = colname if row_num == 0 else data[row_num - 1][colname]
                    _table_text_css(run, text_prop, row_type, txt)
                    if color_grad and isinstance(txt, (int, float)):
                        cell_prop = tbl_style.get('cell', {})
                        gradient = _color.gradient(txt, color_grad)
                        _table_cell_css(row, cell, cell_prop, row_type, gradient)


def oval(shape, spec, data):
    """Function to animate oval shape type."""
    fill = shape.fill
    fill.solid()
    data['_color'] = _color
    bg_color = template(spec['oval']['background-color'], data)
    fill.fore_color.rgb = RGBColor.from_string(bg_color.replace('#', ''))


def chart(shape, spec, data):
    """Replacing chart Data."""
    chart_type = None
    if hasattr(shape.chart, 'chart_type'):
        chart_type = '{}'.format(shape.chart.chart_type).split()[0]

    stacked_or_line = {
        'AREA', 'AREA_STACKED', 'AREA_STACKED_100', 'BAR_CLUSTERED',
        'BAR_OF_PIE', 'BAR_STACKED', 'BAR_STACKED_100', 'COLUMN_CLUSTERED',
        'COLUMN_STACKED', 'COLUMN_STACKED_100', 'CONE_BAR_CLUSTERED',
        'CONE_BAR_STACKED', 'CONE_BAR_STACKED_100', 'CONE_COL',
        'CONE_COL_CLUSTERED', 'CONE_COL_STACKED', 'CONE_COL_STACKED_100',
        'CYLINDER_BAR_CLUSTERED', 'CYLINDER_BAR_STACKED',
        'CYLINDER_BAR_STACKED_100', 'CYLINDER_COL', 'CYLINDER_COL_CLUSTERED',
        'CYLINDER_COL_STACKED', 'CYLINDER_COL_STACKED_100', 'LINE',
        'LINE_MARKERS', 'LINE_MARKERS_STACKED', 'LINE_MARKERS_STACKED_100',
        'LINE_STACKED', 'LINE_STACKED_100', 'PYRAMID_BAR_CLUSTERED',
        'PYRAMID_BAR_STACKED', 'PYRAMID_BAR_STACKED_100', 'PYRAMID_COL',
        'PYRAMID_COL_CLUSTERED', 'PYRAMID_COL_STACKED', 'RADAR_MARKERS',
        'PYRAMID_COL_STACKED_100', 'RADAR', 'RADAR_FILLED', 'PIE',
        'PIE_EXPLODED', 'PIE_OF_PIE'}

    xy_charts = {
        'XY_SCATTER', 'XY_SCATTER_LINES', 'XY_SCATTER_LINES_NO_MARKERS',
        'XY_SCATTER_SMOOTH', 'XY_SCATTER_SMOOTH_NO_MARKERS'}

    bubble_charts = {'BUBBLE', 'BUBBLE_THREE_D_EFFECT'}

    if not chart_type:
        raise NotImplementedError()

    info = spec['chart']
    # Load data
    handler = data.pop('handler') if 'handler' in data else None
    if isinstance(info['x'], (dict,)):
        if 'function' not in info['x']:
            info['x']['function'] = info['x']
        info['x'] = compile_function(info, 'x', data, handler)

    if 'size' in info and isinstance(info['size'], (dict,)):
        if 'function' in info['size']:
            info['size']['function'] = info['size']
        info['size'] = compile_function(info, 'size', data, handler)
    data = compile_function(info, 'data', data, handler)
    change_data = _prepare_data(info, data)
    series_cols = [x for x in change_data.columns if x != info['x']]
    chart_name = None
    # If chart type is stacked bar or line.
    if chart_type in stacked_or_line:
        # Initializing chart data
        chart_name = 'line'
        chart_data = ChartData()
        chart_data.categories = change_data[info['x']].dropna().unique().tolist()
        change_data = _update_chart(info, change_data, chart_data, series_cols)

    # If chart type is scatter plot.
    elif chart_type in xy_charts:
        # Initializing chart data
        chart_name = 'scatter'
        chart_data = XyChartData()
        change_data = _update_chart(info, change_data, chart_data, series_cols, chart='scatter')

    # If chart type is bubble chart.
    elif chart_type in bubble_charts:
        # Initializing chart data
        chart_name = 'bubble'
        chart_data = BubbleChartData()
        change_data = _update_chart(info, change_data, chart_data, series_cols, chart='bubble')
    else:
        raise NotImplementedError()
    shape.chart.replace_data(chart_data)
    is_color = spec['chart'].get('color')

    if isinstance(is_color, (dict,)) and 'function' in is_color:
        is_color = compile_function(spec['chart'], 'color', data, handler)

    if is_color:
        colors = data[spec['chart']['color']].dropna().unique().tolist()
        for index, x in enumerate(shape.chart.series):
            if chart_name in {'scatter', 'bubble'}:
                fill_graph = colors[index].rsplit('#')[-1].lower()
                fill_graph = fill_graph + ('0' * (6 - len(fill_graph)))
                fill = x.marker.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(fill_graph)
                x.marker.format.line.color.rgb = RGBColor.from_string(fill_graph)
            elif chart_name == 'line':
                for i, point in enumerate(x.points):
                    fill_graph = colors[i].rsplit('#')[-1].lower()
                    fill_graph = fill_graph + ('0' * (6 - len(fill_graph)))
                    fill = point.format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(fill_graph)

# Custom Charts Functions below(Sankey, Treemap, Calendarmap).


def sankey(shape, spec, data):
    """Draw sankey in Treemap."""
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    # Getting parent shapes
    pxl_to_inch = 10000
    default_thickness = 40
    spec = spec['sankey']

    handler = data.pop('handler') if 'handler' in data else None
    shapes = shape._parent
    y0 = shape.top
    x0 = shape.left

    width = shape.width
    shape_ids = {'shape': 0}
    height = shape.height

    groups = compile_function(spec, 'groups', data, handler)
    thickness = spec.get('thickness', default_thickness) * pxl_to_inch

    h = (height - (thickness * len(groups))) / (len(groups) - 1) + thickness
    frames = {}
    # Sankey Rectangles and texts.
    sankey_conf = {}
    sankey_conf['x0'] = x0
    sankey_conf['size'] = compile_function(spec, 'size', data, handler)
    sankey_conf['width'] = width
    sankey_conf['order'] = compile_function(spec, 'order', data, handler)
    sankey_conf['text'] = compile_function(spec, 'text', data, handler)
    sankey_conf['color'] = compile_function(spec, 'color', data, handler)
    sankey_conf['attrs'] = spec.get('attrs', {})
    sankey_conf['sort'] = spec.get('sort')
    stroke = spec.get('stroke', '#ffffff')
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    elem_schema = utils.make_element()
    data = compile_function(spec, 'data', data, handler)
    for ibar, group in enumerate(groups):
        y = y0 + h * ibar
        sankey_conf['group'] = [group]
        df = frames[group] = utils.draw_sankey(data, sankey_conf)
        # Adding rectangle
        for key, row in df.iterrows():
            shp = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, row['x'], y, row['width'], thickness)
            rectstyle = {"fill": row['fill'], 'stroke': stroke}
            rect_css(shp, **rectstyle)
            txt_fill = _color.contrast(row['fill'])
            add_text_to_shape(shp, row['text'], spec.get('font-size', 18), txt_fill)

    # Sankey Connection Arcs.
    for ibar, (group1, group2) in enumerate(zip(groups[:-1], groups[1:])):
        sankey_conf['group'] = [group1, group2]
        sankey_conf['sort'] = False
        df = utils.draw_sankey(data, sankey_conf)
        pos = collections.defaultdict(float)
        for key1, row1 in frames[group1].iterrows():
            for key2, row2 in frames[group2].iterrows():
                if (key1, key2) in df.index:
                    row = df.ix[(key1, key2)]
                    y1, y2 = y0 + h * ibar + thickness, y0 + h * (ibar + 1)
                    ym = (y1 + y2) / 2
                    x1 = row1['x'] + pos[0, key1]
                    x2 = row2['x'] + pos[1, key2]

                    _id = shape_ids['shape'] = shape_ids['shape'] + 1
                    shp = utils.cust_shape(
                        0, 0, '{:.0f}'.format(row['width']), '{:.0f}'.format(ym), _id)
                    path = elem_schema['a'].path(
                        w='{:.0f}'.format(row['width']), h='{:.0f}'.format(ym))
                    shp.find('.//a:custGeom', namespaces=elem_schema['nsmap']).append(
                        elem_schema['a'].pathLst(path))
                    path.append(
                        elem_schema['a'].moveTo(elem_schema['a'].pt(
                            x='{:.0f}'.format(x1 + row['width']), y='{:.0f}'.format(y1))))

                    path.append(elem_schema['a'].cubicBezTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x1 + row['width']),
                                            y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x2 + row['width']),
                                            y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x2 + row['width']),
                                            y='{:.0f}'.format(y2))))

                    path.append(elem_schema['a'].lnTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(y2))))

                    path.append(elem_schema['a'].cubicBezTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(y1))))

                    path.append(elem_schema['a'].close())
                    shp.spPr.append(elem_schema['a'].solidFill(
                        utils.fill_color(srgbclr=row['fill'])))
                    shapes._spTree.append(shp)
                    pos[0, key1] += row['width']
                    pos[1, key2] += row['width']


def treemap(shape, spec, data):
    """Function to download data as ppt."""
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    shapes = shape._parent
    x0 = shape.left
    y0 = shape.top
    width = shape.width
    height = shape.height
    spec = spec['treemap']
    stroke = spec.get('stroke', '#ffffff')
    # Load data
    handler = data.pop('handler') if 'handler' in data else None
    spec['keys'] = compile_function(spec, 'keys', data, handler)
    spec['values'] = compile_function(spec, 'values', data, handler)
    spec['size'] = compile_function(spec, 'size', data, handler)
    spec['sort'] = compile_function(spec, 'sort', data, handler)
    spec['color'] = compile_function(spec, 'color', data, handler)
    spec['text'] = compile_function(spec, 'text', data, handler)
    spec['data'] = compile_function(spec, 'data', data, handler)
    # Getting rectangle's width and height using `squarified` algorithm.
    treemap_data = utils.SubTreemap(**spec)
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    font_aspect = 14.5
    pixel_inch = 10000
    default_rect_color = '#cccccc'
    for x, y, w, h, (l, v) in treemap_data.draw(width, height):
        if l == 0:
            shp = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + x0, y + y0, w, h)
            rect_color = spec['color'](v) if spec.get('color') else default_rect_color
            text = spec['text'](v) if spec.get('text') else '{}'.format(v[1])
            rectstyle = {"fill": rect_color, 'stroke': stroke}
            rect_css(shp, **rectstyle)
            txt_fill = _color.contrast(rect_color)
            font_size = min(h, w * font_aspect / fontwidth.fontwidth('{}'.format(text)), pd.np.Inf)
            # Adding text inside rectangles
            add_text_to_shape(shp, text, font_size / pixel_inch, txt_fill)


def calendarmap(shape, spec, data):
    """Draw calendar map in PPT."""
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()

    shapes = shape._parent
    spec = spec['calendarmap']
    handler = data.get('handler')
    # Load data
    data = compile_function(spec, 'data', data, handler)
    startdate = compile_function(spec, 'startdate', data, handler)

    pixel_inch = 10000
    size = spec.get('size', None)

    label_top = spec.get('label_top', 0) * pixel_inch
    label_left = spec.get('label_left', 0) * pixel_inch

    width = spec['width'] * pixel_inch
    shape_top = label_top + shape.top
    shape_left = label_left + shape.left
    y0 = width + shape_top
    x0 = width + shape_left

    # Deleting the shape
    shape.element.delete()
    # Style
    default_color = '#ffffff'
    default_line_color = '#787C74'
    default_txt_color = '#000000'
    font_size = spec.get('font-size', 12)
    stroke = spec.get('stroke', '#ffffff')
    fill_rect = spec.get('fill', '#cccccc')
    text_color = spec.get('text-color', '#000000')
    # Treat infinities as nans when calculating scale
    scaledata = pd.Series(data).replace([pd.np.inf, -pd.np.inf], pd.np.nan)
    lo_data = spec.get('lo', scaledata.min())
    range_data = spec.get('hi', scaledata.max()) - lo_data
    gradient = spec.get('gradient', _color.RdYlGn)
    color = spec.get('color', lambda v: _color.gradient(
        (float(v) - lo_data) / range_data, gradient) if not pd.isnull(v) else default_color)

    startweekday = (startdate.weekday() - spec.get('weekstart', 0)) % 7
    # Weekday Mean and format
    weekday_mean = pd.Series(
        [scaledata[(x - startweekday) % 7::7].mean() for x in range(7)])
    weekday_format = spec.get('format', '{:,.%df}' % utils.decimals(weekday_mean.values))
    # Weekly Mean and format
    weekly_mean = pd.Series([scaledata[max(0, x):x + 7].mean()
                             for x in range(-startweekday, len(scaledata) + 7, 7)])
    weekly_format = spec.get('format', '{:,.%df}' % utils.decimals(weekly_mean.values))

    # Scale sizes as square roots from 0 to max (not lowest to max -- these
    # should be an absolute scale)
    sizes = width * utils.scale(
        [v ** .5 for v in size], lo=0) if size is not None else [width] * len(scaledata)

    for i, val in enumerate(data):
        nx = (i + startweekday) // 7
        ny = (i + startweekday) % 7
        d = startdate + datetime.timedelta(days=i)
        fill = '#ffffff'
        if not pd.isnull(val):
            fill = color(val)

        shp = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x0 + (width * nx) + (width - sizes[i]) / 2,
            y0 + (width * ny) + (width - sizes[i]) / 2,
            sizes[i], sizes[i])
        rectstyle = {"fill": fill, 'stroke': stroke}
        rect_css(shp, **rectstyle)
        add_text_to_shape(shp, '%02d' % d.day, font_size, _color.contrast(fill))

        # Draw the boundary lines between months
        if i >= 7 and d.day == 1 and ny > 0:
            vertical_line = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x0 + width * nx, y0 + (width * ny), width, 2 * pixel_inch)
            rectstyle = {"fill": default_line_color, 'stroke': default_line_color}
            rect_css(vertical_line, **rectstyle)

        if i >= 7 and d.day <= 7 and nx > 0:
            horizontal_line = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x0 + (width * nx), y0 + (width * ny), 2 * pixel_inch, width)
            rectstyle = {"fill": default_line_color, 'stroke': default_line_color}
            rect_css(horizontal_line, **rectstyle)

        # Adding weekdays text to the chart (left side)
        if i < 7:
            txt = shapes.add_textbox(
                x0 - (width / 2), y0 + (width * ny) + (width / 2), width, width)
            add_text_to_shape(txt, d.strftime('%a')[0], font_size, default_txt_color)

        # Adding months text to the chart (top)
        if d.day <= 7 and ny == 0:
            txt = shapes.add_textbox(
                x0 + (width * nx), y0 - (width / 2), width, width)
            add_text_to_shape(txt, d.strftime('%b %Y'), font_size, default_txt_color)

    if label_top:
        lo_weekly = spec.get('lo', weekly_mean.min())
        range_weekly = spec.get('hi', weekly_mean.max()) - lo_weekly
        for nx, val in enumerate(weekly_mean):
            if not pd.isnull(val):
                w = label_top * ((val - lo_weekly) / range_weekly)
                px = x0 + (width * nx)
                top_bar = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, px, shape_top - w, width, w)
                rectstyle = {"fill": fill_rect, 'stroke': stroke}
                rect_css(top_bar, **rectstyle)
                top_txt = shapes.add_textbox(px, shape_top - width, width, width)
                add_text_to_shape(
                    top_txt, weekly_format.format(weekly_mean[nx]), font_size, text_color)

    if label_left:
        lo_weekday = spec.get('lo', weekday_mean.min())
        range_weekday = spec.get('hi', weekday_mean.max()) - lo_weekday
        for ny, val in enumerate(weekday_mean):
            if not pd.isnull(val):
                w = label_left * ((val - lo_weekday) / range_weekday)
                bar = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, shape_left - w, y0 + (width * ny), w, width)
                rectstyle = {"fill": fill_rect, 'stroke': stroke}
                rect_css(bar, **rectstyle)
                left_txt = shapes.add_textbox(shape_left - width, y0 + (width * ny), w, width)
                add_text_to_shape(
                    left_txt, weekday_format.format(weekday_mean[ny]), font_size, text_color)


def bullet(shape, spec, data):
    """Function to plot bullet chart."""
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
            raise NotImplementedError()
    spec = spec['bullet']
    orient = spec.get('orient', 'horizontal')
    if orient not in {'horizontal', 'vertical'}:
            raise NotImplementedError()

    x = shape.left
    y = shape.top
    text = spec.get('text')
    height = shape.height if orient == 'horizontal' else shape.width
    width = shape.width if orient == 'horizontal' else shape.height
    handler = data.get('handler')
    for metric in ['poor', 'average', 'good', 'target']:
        if metric in spec:
            spec[metric] = compile_function(spec, metric, data, handler)
        else:
            spec[metric] = np.nan

    if text:
        text = compile_function({'text': spec['text']}, 'text', data, handler)

    # if isinstance(spec['data'], (dict,)) and 'function' in spec['data']:
    spec['data'] = compile_function(spec, 'data', data, handler)
    gradient = spec.get('gradient', 'Greys')
    # if isinstance(gradient, (dict,)) and 'function' in gradient:
    # gradient = compile_function(spec, 'gradient', data, handler)
    shapes = shape._parent
    shape._sp.delete()

    # lo = np.nanmin([data, target, poor, average, good])
    lo = spec.get('lo', 0)
    hi = spec.get('hi', np.nanmax([spec['data'], spec['target'], spec['poor'],
                                  spec['average'], spec['good']]))
    default_font = 18
    gradient = matplotlib.cm.get_cmap(gradient)
    percentage = {'good': 0.125, 'average': 0.25, 'poor': 0.50}
    for index, metric in enumerate(['good', 'average', 'poor']):
        if not np.isnan(spec[metric]):
            scaled = scale_data(spec.get(metric, np.nan), lo, hi, width)
            _width = scaled if orient == 'horizontal' else height
            _hight = height if orient == 'horizontal' else scaled
            yaxis = y if orient == 'horizontal' else y + (width - scaled)
            _rect = rect(shapes, x, yaxis, _width, _hight)
            rectstyle = {'fill': matplotlib.colors.to_hex(gradient(percentage[metric])),
                         'stroke': matplotlib.colors.to_hex(gradient(percentage[metric]))}
            rect_css(_rect, **rectstyle)

    scaled = scale_data(spec['data'], lo, hi, width)
    _width = scaled if orient == 'horizontal' else height / 2.0
    yaxis = y + height / 4.0 if orient == 'horizontal' else y + (width - scaled)
    xaxis = x if orient == 'horizontal' else x + height / 4.0
    _hight = height / 2.0 if orient == 'horizontal' else scaled
    data_rect = rect(shapes, xaxis, yaxis, _width, _hight)
    rectstyle = {'fill': matplotlib.colors.to_hex(gradient(1.0)),
                 'stroke': matplotlib.colors.to_hex(gradient(1.0))}
    rect_css(data_rect, **rectstyle)

    if text:
        data_text = text(spec['data'])
        parent = data_rect._parent
        _factor = 1.25
        text_width_height = _width if orient == 'vertical' else _hight
        _xaxis = x + (_width / 4.0) if orient == 'vertical' else x + scaled - (_hight * _factor)
        parent = parent.add_textbox(_xaxis, yaxis, text_width_height, text_width_height)
        add_text_to_shape(
            parent, data_text, spec.get('fontsize', default_font),
            _color.contrast(matplotlib.colors.to_hex(gradient(1.0))))

    if not np.isnan(spec['target']):
        line_hight = 10000
        scaled = scale_data(spec['target'], lo, hi, width)
        _width = line_hight if orient == 'horizontal' else height
        _hight = height if orient == 'horizontal' else line_hight
        yaxis = y if orient == 'horizontal' else (width - scaled) + y
        xaxis = x + scaled if orient == 'horizontal' else x
        target_line = rect(shapes, xaxis, yaxis, _width, _hight)
        rectstyle = {'fill': matplotlib.colors.to_hex(gradient(1.0)),
                     'stroke': matplotlib.colors.to_hex(gradient(1.0))}
        rect_css(target_line, **rectstyle)
        if text:
            handler = data.get('handler')
            target_text = text(spec['target'])
            parent = target_line._parent
            yaxis = yaxis - (_width / 2) if orient == 'vertical' else yaxis + (_hight / 5.0)
            text_width_height = _width if orient == 'vertical' else _hight
            parent = parent.add_textbox(xaxis, yaxis, text_width_height, text_width_height)
            add_text_to_shape(
                parent, target_text, spec.get('fontsize', default_font),
                _color.contrast(matplotlib.colors.to_hex(gradient(percentage['good']))))


def rectangle(shape, spec, data):
    """Function to modify a rectangle's property in PPT."""
    pxl_to_inch = 10000
    spec = spec['rectangle']

    handler = data.pop('handler') if 'handler' in data else None
    data = compile_function(spec, 'data', data, handler)

    for prop in {'width', 'height', 'top', 'left'}:
        setprop = spec.get(prop)
        if setprop:
            if not isinstance(spec[prop], (dict,)):
                spec[prop] = {'function': '{}'.format(spec[prop]) if not isinstance(
                              spec[prop], (str, six.string_types,)) else spec[prop]}
            setprop = compile_function(spec, prop, data, handler)
            setprop = setprop * pxl_to_inch
        else:
            setprop = getattr(shape, prop)
        setattr(shape, prop, setprop)

    for style in {'fill', 'stroke'}:
        _style = spec.get(style)
        if _style:
            if not isinstance(_style, (dict,)):
                spec[style] = {'function': _style}
            spec[style] = compile_function(spec, style, data, handler)
            if callable(spec[style]):
                spec[style] = spec[style](spec[style])
    rect_css(shape, **spec)


cmdlist = {
    'text': text,
    'oval': oval,
    'image': image,
    'chart': chart,
    'table': table,
    'sankey': sankey,
    'bullet': bullet,
    'treemap': treemap,
    'replace': replace,
    'rectangle': rectangle,
    'calendarmap': calendarmap
}
